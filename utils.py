import requests
from langchain_community.llms import Ollama
import ollama
import fitz  # PyMuPDF (for PDFs)
from docx import Document  # For DOCX files
import pandas as pd  # For Excel files
import io
from pptx import Presentation  # For PPTX files
import ffmpeg
import os
import whisper
import tempfile
import json
import bcrypt


USERS_FILE = "users.json"


def extract_text_from_pptx(file):
    byte_data = file.read()
    prs = Presentation(io.BytesIO(byte_data))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_pdf(uploaded_file):
    byte_data = uploaded_file.read()
    with io.BytesIO(byte_data) as byte_file:
        doc = fitz.open(stream=byte_file, filetype="pdf")
        text = ""
        for page_num in range(doc.page_count):
            page = doc.load_page(page_num)
            text += page.get_text()
    return text

def extract_text_from_word(file):
    byte_data = file.read()
    doc = Document(io.BytesIO(byte_data))
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def extract_text_from_excel(file):
    byte_data = file.read()
    df = pd.read_excel(io.BytesIO(byte_data))
    text = df.to_string()
    return text

def extract_audio_from_mp4(video_path, output_audio_path):
    ffmpeg.input(video_path).output(output_audio_path).run(overwrite_output=True)
    os.remove(video_path)

def transcribe_audio(audio_path):
    model = whisper.load_model("base")
    result = model.transcribe(audio_path)
    return result['text']

def extract_text_from_file(file):
    if file.type == "application/pdf":
        return extract_text_from_pdf(file)
    elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return extract_text_from_word(file)
    elif file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        return extract_text_from_excel(file)
    elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return extract_text_from_pptx(file)
    elif file.type == "video/mp4":
        with tempfile.NamedTemporaryFile(delete=False, suffix=".mp4") as tmp_video:
            tmp_video.write(file.read())
            tmp_video_path = tmp_video.name

        with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as temp_audio:
            temp_audio_path = temp_audio.name

        extract_audio_from_mp4(tmp_video_path, temp_audio_path)
        return transcribe_audio(temp_audio_path)
    else:
        return "Unsupported file type"

model = "llama3.1:8b-instruct-q4_K_M"

system_prompt =  """
You are an expert in Battery Management Systems (BMS), specializing in making complex information easy to understand. 
When provided with a question and relevant context about Battery Management Systems, your task is to deliver a clear, concise, and accurate response strictly based on the provided context.
Your objective is to ensure the user fully understands the key points and concepts related to BMS, without introducing any additional information that is not contained in the context.
Always answer in the language of the question provided, not the language of the context!

answer the following question based on the facts present in the context provided:

Context: ```{context}```

question: ```{question}```

Provide your answer **directly** based on the facts in the context, without any interpretation or creative generation. If the context does not contain relevant information, explicitly state that you cannot answer the question.
"""

def talk_to_llama3(actual_prompt, temp):
    print('^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^')
    print(actual_prompt)
    print('^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^^')
    
    response = ollama.chat(
        model=model, 
        messages=[{'role': 'system', 'content': system_prompt},
                  {'role': 'user', 'content': actual_prompt}]
    )
    return response['message']['content']

def query_rag(prompt):
    url = f'http://localhost:8000/chat?question={prompt}'
    response = requests.get(url)
    data = response.json()
    return data.get('reply', "❌ Erreur lors de la récupération de la réponse.")
    

def query_rag_with_context(prompt, context):
    final_prompt = f"""Context: {context}

Question: {prompt}"""

    return talk_to_llama3(final_prompt, temp=0.7)

def remove_duplicates(l):
    l_without_duplicates = []
    unseen = set()
    for element in l:
        if element[0] not in unseen:
            l_without_duplicates.append(element)
            unseen.add(element[0])
    return l_without_duplicates

def get_similiar(queries, db, n=5):
    results = []
    for query in queries:
        result = db.max_marginal_relevance_search(query, k=n)
        results.extend(result)
    results = [(doc.page_content, 1) for doc in results]
    return results

def reformulate(question, REFORMULATION_TEMP):
    questions = [question]
    for i in range(1):
        q = f"""
        Réformulez cette question,
        Veuillez le reformuler en tenant compte de ce contexte et en ajoutant quelques mots clés en lien avec cela.
        Assurez-vous d'utiliser des mots différents de ceux de la question posée.
        votre réponse doit se composer uniquement des questions formulées et rien d'autre.
        question: {question}
        reformulation:
        """
        questions.append(talk_to_llama3(q, REFORMULATION_TEMP))
    return questions

#------------------------SIGNUP AND LOGIN ---------------------------

def load_users():
    try:
        with open(USERS_FILE, "r") as f:
            return json.load(f)
    except (FileNotFoundError, json.JSONDecodeError):
        return {}

def save_users(users):
    with open(USERS_FILE, "w") as f:
        json.dump(users, f, indent=4)

def hash_password(password):
    return bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')

def verify_password(password, hashed):
    return bcrypt.checkpw(password.encode('utf-8'), hashed.encode('utf-8'))

def add_user(username, password):
    users = load_users()
    if username in users:
        return False  # User already exists
    users[username] = hash_password(password)
    save_users(users)
    return True

def authenticate_user(username, password):
    users = load_users()
    if username in users and verify_password(password, users[username]):
        return True
    return False

#------------------------------- FEEDBACK ----------------------
# Path to store feedback
FEEDBACK_FILE = "feedback_data.json"

# Model configuration
model = "llama3.1:8b-instruct-q4_K_M"
system_prompt = """
You are an expert in Battery Management Systems (BMS), specializing in making complex information easy to understand. 
When provided with a question and relevant context about Battery Management Systems, your task is to deliver a clear, concise, and accurate response strictly based on the provided context.
Your objective is to ensure the user fully understands the key points and concepts related to BMS, without introducing any additional information that is not contained in the context.
Always answer in the language of the question provided, not the language of the context!
"""

# Function to store feedback data (including thumbs up/thumbs down)
def store_feedback(feedback_data):
    """
    Store feedback data in a JSON file.

    Args:
    - feedback_data (dict): Contains the feedback type, response, and user question.
    """
    # Open the file and append feedback
    if os.path.exists(FEEDBACK_FILE):
        with open(FEEDBACK_FILE, "a") as f:
            json.dump(feedback_data, f, indent=4)
            f.write("\n")  # To separate each feedback entry
    else:
        # Create file and write feedback
        with open(FEEDBACK_FILE, "w") as f:
            json.dump([feedback_data], f, indent=4)

# Function to regenerate the answer (when thumbs down is given)
def regenerate_answer(question, temp=0.7):
    """
    Regenerate the answer when the thumbs down feedback is received.

    Args:
    - question (str): The user's original question.
    - temp (float): Temperature parameter for the model (controls creativity of the answer).
    
    Returns:
    - (str): The new response generated by the model.
    """
    actual_prompt = f"""
    Context: {question}
    Question: {question}

    Provide a new response based on the user's query.
    """
    
    response = ollama.chat(
        model=model,
        messages=[{'role': 'system', 'content': system_prompt},
                  {'role': 'user', 'content': actual_prompt}]
    )
    return response['message']['content']